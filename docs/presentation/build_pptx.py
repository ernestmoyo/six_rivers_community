"""
Six Rivers Community Intelligence Platform — Methodology Presentation
Clean, readable, white background, native PowerPoint shapes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colours ──
NAVY      = RGBColor(0x1E, 0x3A, 0x5F)
GREEN     = RGBColor(0x2E, 0x7D, 0x32)
TEAL      = RGBColor(0x00, 0x79, 0x6B)
ORANGE    = RGBColor(0xE6, 0x51, 0x00)
RED       = RGBColor(0xC6, 0x28, 0x28)
BLUE      = RGBColor(0x15, 0x65, 0xC0)
GRAY      = RGBColor(0x75, 0x75, 0x75)
DARK      = RGBColor(0x21, 0x21, 0x21)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xFA, 0xFA, 0xFA)
CARD_BG   = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_BLUE = RGBColor(0xBB, 0xDE, 0xFB)
LIGHT_GREEN = RGBColor(0xC8, 0xE6, 0xC9)
LIGHT_ORANGE = RGBColor(0xFF, 0xE0, 0xB2)
LIGHT_RED = RGBColor(0xFF, 0xCD, 0xD2)
LIGHT_TEAL = RGBColor(0xB2, 0xDF, 0xDB)

IMG_DIR = os.path.join(os.path.dirname(__file__), "images")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──

def white_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

def bar(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def rounded_box(slide, x, y, w, h, fill_color, border_color=None, border_width=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = border_width
    else:
        s.line.fill.background()
    return s

def txt(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def multi_txt(slide, x, y, w, h, lines, size=18, color=DARK, line_spacing=Pt(10), bold_first=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = line_spacing
        if bold_first and i == 0:
            p.font.bold = True
    return tf

def title_bar(slide, title, subtitle=None):
    bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), NAVY)
    txt(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.7), title,
        size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.45), subtitle,
            size=18, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)

def footer(slide):
    bar(slide, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45), NAVY)
    txt(slide, Inches(0.5), Inches(7.07), Inches(12.3), Inches(0.4),
        "Six Rivers Community Intelligence Platform   |   7Square Inc.   |   March 2026",
        size=11, color=WHITE, align=PP_ALIGN.CENTER)

def section_divider(slide, number, title_text, subtitle_text):
    """Full-slide section divider — dark navy with large number."""
    white_bg(slide)
    bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    # Large faded number
    txt(slide, Inches(0.8), Inches(1.0), Inches(4), Inches(2.5), number,
        size=120, bold=True, color=RGBColor(0x2A, 0x4A, 0x70), align=PP_ALIGN.LEFT)
    # Title
    txt(slide, Inches(0.8), Inches(3.2), Inches(11), Inches(1.2), title_text,
        size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    # Green accent line
    bar(slide, Inches(0.8), Inches(4.5), Inches(3), Inches(0.06), GREEN)
    # Subtitle
    txt(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.8), subtitle_text,
        size=20, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(s)
rounded_box(s, Inches(1.2), Inches(0.8), Inches(10.9), Inches(5.9), NAVY)
bar(s, Inches(5.5), Inches(4.1), Inches(2.3), Inches(0.05), GREEN)

txt(s, Inches(1.8), Inches(1.6), Inches(9.7), Inches(1.0),
    "COMMUNITY INTELLIGENCE", 48, True, WHITE, PP_ALIGN.CENTER)
txt(s, Inches(1.8), Inches(2.5), Inches(9.7), Inches(0.8),
    "PLATFORM", 48, True, WHITE, PP_ALIGN.CENTER)
txt(s, Inches(1.8), Inches(3.4), Inches(9.7), Inches(0.6),
    "Methodology & Approach", 26, False, LIGHT_BLUE, PP_ALIGN.CENTER)
txt(s, Inches(1.8), Inches(4.5), Inches(9.7), Inches(0.5),
    "Six Rivers Africa  +  7Square Inc.", 22, False, WHITE, PP_ALIGN.CENTER)
txt(s, Inches(1.8), Inches(5.3), Inches(9.7), Inches(0.5),
    "March 2026", 18, False, RGBColor(0x90, 0xCA, 0xF9), PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM (section divider)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "01", "The Challenge", "Why does the community team need its own tool?")


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — THE CHALLENGE (content)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(s)
title_bar(s, "The Challenge")
footer(s)

# Left column — 3 challenge cards
challenges_left = [
    ("No digital tools for community work", "Field officers use notebooks, WhatsApp, and scattered spreadsheets. No centralised view of farming activities or wildlife incidents.", LIGHT_BLUE, BLUE),
    ("Elephant crop raids", "Villages adjacent to Nyerere National Park face frequent elephant intrusions that destroy crops and threaten livelihoods.", LIGHT_RED, RED),
    ("Cattle pressure in Mbarali", "Growing challenge near Ruaha National Park. Needs systematic monitoring and evidence for policy.", LIGHT_ORANGE, ORANGE),
]

for i, (title, desc, bg_col, accent) in enumerate(challenges_left):
    y = Inches(1.7 + i * 1.75)
    rounded_box(s, Inches(0.6), y, Inches(5.8), Inches(1.5), bg_col, accent, Pt(2))
    txt(s, Inches(1.0), y + Inches(0.15), Inches(5.0), Inches(0.4), title,
        size=20, bold=True, color=accent)
    txt(s, Inches(1.0), y + Inches(0.6), Inches(5.0), Inches(0.8), desc,
        size=15, color=DARK)

# Right column — 3 challenge cards
challenges_right = [
    ("Existing tools look INSIDE parks", "The Landscape Dashboard monitors satellite data inside protected areas. The community team works OUTSIDE — in villages, farms, and nurseries.", LIGHT_GREEN, GREEN),
    ("Donor reporting is manual", "Mary's quarterly reports require manual aggregation across all activities. Time-consuming and error-prone.", LIGHT_BLUE, BLUE),
    ("Dry-season planting failures", "42% seedling loss during dry season. Planning needs seasonal awareness to time interventions correctly.", LIGHT_TEAL, TEAL),
]

for i, (title, desc, bg_col, accent) in enumerate(challenges_right):
    y = Inches(1.7 + i * 1.75)
    rounded_box(s, Inches(6.9), y, Inches(5.8), Inches(1.5), bg_col, accent, Pt(2))
    txt(s, Inches(7.3), y + Inches(0.15), Inches(5.0), Inches(0.4), title,
        size=20, bold=True, color=accent)
    txt(s, Inches(7.3), y + Inches(0.6), Inches(5.0), Inches(0.8), desc,
        size=15, color=DARK)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — OUR APPROACH (section divider)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "02", "Our Approach", "A four-phase methodology built around your team's real needs")


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — METHODOLOGY: Listen → Design → Build → Iterate
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(s)
title_bar(s, "Methodology", "Listen  —  Design  —  Build  —  Iterate")
footer(s)

phases = [
    ("01", "LISTEN", GREEN, LIGHT_GREEN, [
        "Met with field team",
        "(Lilian, Justina, Irene,",
        " Mary, Edna, Thomas)",
        "",
        "Understood real operations:",
        "  Work OUTSIDE parks",
        "  Small farms (< 0.4 ha)",
        "  3-month crop cycles",
        "",
        "Mapped 21 real villages",
        "across 2 district councils",
    ]),
    ("02", "DESIGN", BLUE, LIGHT_BLUE, [
        "5 farming approach modules",
        "tailored to your activities",
        "",
        "2 operational zones:",
        "  Ifakara Town Council",
        "  Mbarali District Council",
        "",
        "Wildlife incident tracking",
        "Seasonal failure awareness",
        "Impact reporting framework",
        "for donor requirements",
    ]),
    ("03", "BUILD", ORANGE, LIGHT_ORANGE, [
        "Next.js web application",
        "(works on phone + desktop)",
        "",
        "Mapbox satellite maps",
        "with real village locations",
        "",
        "Dashboard with live KPIs",
        "Farming approaches module",
        "Nursery & seedling tracking",
        "Cattle pressure monitoring",
        "Donor report generation",
    ]),
    ("04", "ITERATE", RED, LIGHT_RED, [
        "Present demo to team",
        "Gather field feedback",
        "",
        "Priority next steps:",
        "  Database + authentication",
        "  Mobile data entry forms",
        "  GPS for fences & plots",
        "  Offline mode for field",
        "  Swahili language support",
        "",
        "Deploy to production",
    ]),
]

for i, (num, label, accent, bg_col, items) in enumerate(phases):
    x = Inches(0.4 + i * 3.2)
    # Card background
    rounded_box(s, x, Inches(1.65), Inches(2.95), Inches(5.1), bg_col, accent, Pt(2.5))
    # Phase number + label
    txt(s, x + Inches(0.2), Inches(1.8), Inches(2.5), Inches(0.5), num,
        size=32, bold=True, color=accent, align=PP_ALIGN.LEFT)
    txt(s, x + Inches(1.0), Inches(1.85), Inches(1.8), Inches(0.4), label,
        size=22, bold=True, color=accent, align=PP_ALIGN.LEFT)
    # Accent line
    bar(s, x + Inches(0.2), Inches(2.35), Inches(2.5), Inches(0.04), accent)
    # Content lines
    multi_txt(s, x + Inches(0.2), Inches(2.55), Inches(2.5), Inches(4.0),
              items, size=14, color=DARK, line_spacing=Pt(3))

# Arrow connectors between cards
for i in range(3):
    x = Inches(3.35 + i * 3.2)
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(4.0), Inches(0.35), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY
    arrow.line.fill.background()


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — OPERATIONAL ZONES (section divider)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "03", "Operational Zones", "21 villages across 2 district councils adjacent to national parks")


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — OPERATIONAL ZONES (content)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(s)
title_bar(s, "Operational Zones", "21 villages across Ifakara Town Council and Mbarali District Council")
footer(s)

# LEFT — Ifakara TC
rounded_box(s, Inches(0.5), Inches(1.65), Inches(6.0), Inches(5.1), LIGHT_GREEN, GREEN, Pt(2.5))
txt(s, Inches(0.8), Inches(1.8), Inches(5.4), Inches(0.5),
    "IFAKARA TOWN COUNCIL", 26, True, GREEN)
txt(s, Inches(0.8), Inches(2.3), Inches(5.4), Inches(0.4),
    "13 villages  |  Adjacent to Nyerere National Park", 16, False, GRAY)
bar(s, Inches(0.8), Inches(2.75), Inches(5.4), Inches(0.03), GREEN)

ifakara_villages = [
    "Kiberege Ward:  Nyamwezi,  Mkasu,  Bwawani",
    "Kibaoni Ward:  Lugongole",
    "Mwaya Ward:  Mhelule,  Mikoleko",
    "Sanje Ward:  Miwangani",
    "Kisawasawa Ward:  Mpanga",
    "Mang'ula A Ward:  Msalise",
    "Signal Ward:  Sagamaganga,  Signal",
    "Other:  Katindiuka,  Mbasa",
]
multi_txt(s, Inches(0.8), Inches(2.95), Inches(5.4), Inches(3.5),
          ifakara_villages, size=16, color=DARK, line_spacing=Pt(8))

# RIGHT — Mbarali DC
rounded_box(s, Inches(6.85), Inches(1.65), Inches(6.0), Inches(5.1), LIGHT_ORANGE, ORANGE, Pt(2.5))
txt(s, Inches(7.15), Inches(1.8), Inches(5.4), Inches(0.5),
    "MBARALI DISTRICT COUNCIL", 26, True, ORANGE)
txt(s, Inches(7.15), Inches(2.3), Inches(5.4), Inches(0.4),
    "8 villages  |  Adjacent to Ruaha National Park", 16, False, GRAY)
bar(s, Inches(7.15), Inches(2.75), Inches(5.4), Inches(0.03), ORANGE)

mbarali_villages = [
    "Miyombweni Ward:  Magigiwe",
    "Kiberege Ward:  Mapogoro,  Mlungu",
    "Madibira Ward:  Iheha",
    "Mahango Ward:  Chalisuka",
    "Mkunywa Ward:  Mwaya",
    "Other:  Ikoga Mpya,  Nyakadete",
    "",
    "Key issue: Cattle pressure near Ruaha NP",
]
multi_txt(s, Inches(7.15), Inches(2.95), Inches(5.4), Inches(3.5),
          mbarali_villages, size=16, color=DARK, line_spacing=Pt(8))


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — FARMING APPROACHES (section divider)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "04", "Farming Approaches", "Five integrated modules for human-wildlife coexistence")


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — FARMING APPROACHES (content)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(s)
title_bar(s, "Farming Approaches", "Five modules aligned to your real field activities")
footer(s)

approaches = [
    ("Chilli Fencing", "Nature-based\nelephant deterrence", "Individual farmer\nperimeter planting", "Track fences,\ndeterrence events,\nreplanting needs", RED, LIGHT_RED),
    ("Agroforestry", "Tree-crop\nintegration", "Cocoa, fruit trees\n+ food crops", "Track species mix,\nsurvival rates,\nyields per plot", GREEN, LIGHT_GREEN),
    ("Shambachungu", "Group-based\nwildlife-friendly farming", "Shared plots,\ncollective management", "Track groups,\nmembers, crops,\ntree species", ORANGE, LIGHT_ORANGE),
    ("Horticulture", "Vegetable\nproduction", "Onions, tomatoes,\nleafy greens", "Track varieties,\nharvest yields,\nmarket sales", TEAL, LIGHT_TEAL),
    ("Tree Nurseries", "Community\nseedling production", "Indigenous + exotic\nfor restoration", "Track batches,\nsurvival, distribution,\nseasonal losses", BLUE, LIGHT_BLUE),
]

for i, (name, desc, detail, tracking, accent, bg_col) in enumerate(approaches):
    x = Inches(0.3 + i * 2.6)
    w = Inches(2.35)
    # Card
    rounded_box(s, x, Inches(1.65), w, Inches(5.1), bg_col, accent, Pt(2.5))
    # Title
    txt(s, x + Inches(0.15), Inches(1.8), w - Inches(0.3), Inches(0.45),
        name, 20, True, accent, PP_ALIGN.CENTER)
    bar(s, x + Inches(0.15), Inches(2.3), w - Inches(0.3), Inches(0.03), accent)
    # Description
    txt(s, x + Inches(0.15), Inches(2.5), w - Inches(0.3), Inches(0.8),
        desc, 14, False, DARK, PP_ALIGN.CENTER)
    # Detail
    txt(s, x + Inches(0.15), Inches(3.5), w - Inches(0.3), Inches(0.35),
        "What it looks like:", 12, True, accent, PP_ALIGN.LEFT)
    txt(s, x + Inches(0.15), Inches(3.85), w - Inches(0.3), Inches(0.8),
        detail, 14, False, DARK, PP_ALIGN.LEFT)
    # Tracking
    txt(s, x + Inches(0.15), Inches(4.8), w - Inches(0.3), Inches(0.35),
        "What we track:", 12, True, accent, PP_ALIGN.LEFT)
    txt(s, x + Inches(0.15), Inches(5.15), w - Inches(0.3), Inches(1.2),
        tracking, 14, False, DARK, PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — PLATFORM STRUCTURE (section divider)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "05", "Platform Structure", "How the solution is built — from field data to decisions")


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — PLATFORM ARCHITECTURE (native shapes)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(s)
title_bar(s, "Platform Architecture", "Four layers: Field  >  Collection  >  Intelligence  >  Outputs")
footer(s)

layers = [
    ("FIELD LEVEL", GREEN, LIGHT_GREEN,
     [("Field Officers", "Lilian, Justina, Irene"),
      ("Farmers", "145 across 21 villages"),
      ("Village Leaders", "Community liaison")]),
    ("DATA COLLECTION", ORANGE, LIGHT_ORANGE,
     [("Farm Visits", "Seedling checks, plot surveys"),
      ("Fence Inspections", "Chilli fence status"),
      ("Wildlife Reports", "Incident documentation"),
      ("Nursery Monitoring", "Batch tracking")]),
    ("INTELLIGENCE PLATFORM", BLUE, LIGHT_BLUE,
     [("Database", "PostgreSQL + PostGIS"),
      ("Interactive Maps", "Mapbox GL JS satellite"),
      ("Analytics", "KPIs, trends, comparisons")]),
    ("OUTPUTS", RED, LIGHT_RED,
     [("Live Dashboard", "Real-time field view"),
      ("Donor Reports", "Quarterly impact PDFs"),
      ("Alerts", "Early warnings")]),
]

for row, (layer_name, accent, bg_col, items) in enumerate(layers):
    y = Inches(1.6 + row * 1.35)
    # Layer label
    rounded_box(s, Inches(0.4), y, Inches(2.4), Inches(1.1), accent)
    txt(s, Inches(0.5), y + Inches(0.25), Inches(2.2), Inches(0.6),
        layer_name, 16, True, WHITE, PP_ALIGN.CENTER)

    # Items in that layer
    item_width = Inches(2.3)
    start_x = Inches(3.2)
    for i, (item_name, item_desc) in enumerate(items):
        x = start_x + Inches(i * 2.55)
        rounded_box(s, x, y, item_width, Inches(1.1), bg_col, accent, Pt(1.5))
        txt(s, x + Inches(0.1), y + Inches(0.1), item_width - Inches(0.2), Inches(0.4),
            item_name, 15, True, accent, PP_ALIGN.CENTER)
        txt(s, x + Inches(0.1), y + Inches(0.5), item_width - Inches(0.2), Inches(0.5),
            item_desc, 12, False, GRAY, PP_ALIGN.CENTER)

    # Down arrow (except last row)
    if row < 3:
        arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
            Inches(1.35), y + Inches(1.1), Inches(0.3), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — DATA FLOW
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(s)
title_bar(s, "Data Flow", "How information moves from the field to actionable decisions")
footer(s)

img_path = os.path.join(IMG_DIR, "dataflow.png")
if os.path.exists(img_path):
    slide_w = Inches(12.0)
    slide_h = Inches(4.8)
    s.shapes.add_picture(img_path, Inches(0.65), Inches(1.7), slide_w, slide_h)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — IMPACT MEASUREMENT (section divider)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "06", "Impact Measurement", "Honest numbers from the ground — what we track and why it matters")


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — CURRENT NUMBERS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(s)
title_bar(s, "Current State — Pilot Data", "Real numbers, honestly reported")
footer(s)

kpis_row1 = [
    ("145", "Farmers Registered", GREEN),
    ("21", "Villages Covered", BLUE),
    ("2,400", "Seedlings Distributed", ORANGE),
    ("58%", "Survival Rate", RED),
]

kpis_row2 = [
    ("18", "Active Chilli Fences", RED),
    ("78.5%", "Deterrence Success", GREEN),
    ("6", "Shambachungu Groups", ORANGE),
    ("3", "Community Nurseries", BLUE),
]

for row_idx, kpis in enumerate([kpis_row1, kpis_row2]):
    y = Inches(1.7 + row_idx * 2.65)
    for i, (value, label, accent) in enumerate(kpis):
        x = Inches(0.5 + i * 3.15)
        rounded_box(s, x, y, Inches(2.85), Inches(2.3), CARD_BG, accent, Pt(3))
        txt(s, x, y + Inches(0.3), Inches(2.85), Inches(1.0),
            value, 52, True, accent, PP_ALIGN.CENTER)
        txt(s, x, y + Inches(1.4), Inches(2.85), Inches(0.6),
            label, 18, False, DARK, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 15 — WHO BENEFITS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(s)
title_bar(s, "Who Benefits", "Every stakeholder gets what they need")
footer(s)

beneficiaries = [
    ("Field Officers", "Real-time operational dashboard\nMobile data entry\nGPS-tagged records\nOffline capability", GREEN, LIGHT_GREEN),
    ("Impact Manager", "Automated donor reports\nQuarterly KPI summaries\nTrend analysis over time\nEvidence for proposals", BLUE, LIGHT_BLUE),
    ("Communities", "Reduced crop raids\nBetter farming outcomes\nGroup coordination\nAccess to seedlings", ORANGE, LIGHT_ORANGE),
    ("Management", "Strategic overview\nCross-zone comparison\nResource allocation\nEarly warning system", RED, LIGHT_RED),
]

for i, (who, what, accent, bg_col) in enumerate(beneficiaries):
    x = Inches(0.4 + i * 3.2)
    rounded_box(s, x, Inches(1.65), Inches(2.95), Inches(5.1), bg_col, accent, Pt(2.5))
    txt(s, x + Inches(0.15), Inches(1.85), Inches(2.65), Inches(0.5),
        who, 24, True, accent, PP_ALIGN.CENTER)
    bar(s, x + Inches(0.15), Inches(2.4), Inches(2.65), Inches(0.03), accent)
    multi_txt(s, x + Inches(0.25), Inches(2.65), Inches(2.45), Inches(3.8),
              what.split("\n"), size=17, color=DARK, line_spacing=Pt(12))


# ════════════════════════════════════════════════════════════════
# SLIDE 16 — WHAT'S NEXT (section divider)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(s, "07", "What's Next", "Roadmap for the next iteration")


# ════════════════════════════════════════════════════════════════
# SLIDE 17 — ROADMAP
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(s)
title_bar(s, "Roadmap", "Immediate priorities and short-term goals")
footer(s)

# LEFT — Immediate
rounded_box(s, Inches(0.5), Inches(1.65), Inches(5.9), Inches(5.1), LIGHT_GREEN, GREEN, Pt(2.5))
txt(s, Inches(0.8), Inches(1.85), Inches(5.3), Inches(0.5),
    "IMMEDIATE  (Next 2 weeks)", 24, True, GREEN)
bar(s, Inches(0.8), Inches(2.4), Inches(5.3), Inches(0.03), GREEN)

immediate = [
    "Connect to PostgreSQL + PostGIS database",
    "User authentication (field officers, managers)",
    "Mobile-friendly data entry forms",
    "GPS capture for chilli fences and farm plots",
    "Deploy live version to Vercel",
]
multi_txt(s, Inches(0.8), Inches(2.6), Inches(5.3), Inches(3.8),
          immediate, size=18, color=DARK, line_spacing=Pt(14))

# RIGHT — Short-term
rounded_box(s, Inches(6.9), Inches(1.65), Inches(5.9), Inches(5.1), LIGHT_BLUE, BLUE, Pt(2.5))
txt(s, Inches(7.2), Inches(1.85), Inches(5.3), Inches(0.5),
    "SHORT-TERM  (1–2 months)", 24, True, BLUE)
bar(s, Inches(7.2), Inches(2.4), Inches(5.3), Inches(0.03), BLUE)

shortterm = [
    "Automated quarterly donor PDF reports",
    "Swahili language support",
    "Full offline mode for field use",
    "Photo uploads for wildlife incidents",
    "Integration with Landscape Dashboard",
]
multi_txt(s, Inches(7.2), Inches(2.6), Inches(5.3), Inches(3.8),
          shortterm, size=18, color=DARK, line_spacing=Pt(14))


# ════════════════════════════════════════════════════════════════
# SLIDE 18 — QUESTIONS / TRANSITION TO DEMO
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(s)
rounded_box(s, Inches(1.2), Inches(0.8), Inches(10.9), Inches(5.9), NAVY)
bar(s, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.05), GREEN)

txt(s, Inches(1.8), Inches(2.0), Inches(9.7), Inches(1.0),
    "Questions & Discussion", 44, True, WHITE, PP_ALIGN.CENTER)
txt(s, Inches(1.8), Inches(3.2), Inches(9.7), Inches(0.6),
    "followed by", 20, False, LIGHT_BLUE, PP_ALIGN.CENTER)
txt(s, Inches(1.8), Inches(3.8), Inches(9.7), Inches(0.6),
    "Live Demo of the Platform", 32, True, WHITE, PP_ALIGN.CENTER)
txt(s, Inches(1.8), Inches(5.0), Inches(9.7), Inches(0.5),
    "Ernest Moyo   |   ernest@7squareinc.com   |   7Square Inc.", 18, False,
    RGBColor(0x90, 0xCA, 0xF9), PP_ALIGN.CENTER)


# ── Save ──
output_path = os.path.join(os.path.dirname(__file__), "Six_Rivers_Methodology_Presentation.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
